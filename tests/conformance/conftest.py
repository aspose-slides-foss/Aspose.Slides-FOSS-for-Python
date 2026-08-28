"""Fixtures for the package-level conformance tests."""

from __future__ import annotations

import os

import pytest

from .harness import ProducedPackage, verify_library_not_shadowed

verify_library_not_shadowed()

from aspose.slides_foss import Presentation, ShapeType  # noqa: E402
from aspose.slides_foss.export import SaveFormat  # noqa: E402


def _require_python_pptx():
    pytest.importorskip(
        "pptx",
        reason="python-pptx provides the third-party read-back; install the "
               "test extra: pip install -e '.[test]'",
    )


@pytest.fixture()
def produced(tmp_path):
    """Save a presentation and hand back the file, opened for inspection.

    Usage::

        pkg = produced(pres)                       # saves as .pptx
        pkg = produced(pres, SaveFormat.POTX, "deck.potx")

    The presentation is disposed after saving, so a test cannot accidentally
    fall back on the in-memory object and assert against the library's own idea
    of what it wrote.
    """
    opened: list[ProducedPackage] = []

    def _save(pres, save_format: SaveFormat = SaveFormat.PPTX, name: str = "produced.pptx"):
        path = os.path.join(str(tmp_path), name)
        pres.save(path, save_format)
        pres.dispose()
        package = ProducedPackage(path)
        opened.append(package)
        return package

    yield _save

    # An archive left open holds a handle inside tmp_path, and Windows refuses
    # to remove a directory that has one; the failure then lands on teardown
    # rather than on the test that opened it.
    for package in opened:
        package.close()


@pytest.fixture()
def blank_presentation():
    """A one-slide presentation straight from the bundled template."""
    return Presentation()


@pytest.fixture()
def shape_on_blank_slide():
    """Return ``(presentation, shape)`` — a rectangle on the first slide."""

    def _make(pres=None, with_text: str | None = None):
        pres = pres or Presentation()
        shape = pres.slides[0].shapes.add_auto_shape(
            ShapeType.RECTANGLE, 50.0, 50.0, 200.0, 100.0
        )
        if with_text is not None:
            shape.add_text_frame(with_text)
        return pres, shape

    return _make


@pytest.fixture()
def python_pptx():
    """The third-party reader, skipped with a useful message if absent."""
    _require_python_pptx()
    import pptx

    return pptx
